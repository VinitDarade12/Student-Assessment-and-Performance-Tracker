import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
ARTIFACT_DIR = r"C:\Users\admin\.gemini\antigravity\brain\bf8febe9-f6c7-4d8e-aa76-c4be6d85ce1d"
BG_IMAGE = os.path.join(ARTIFACT_DIR, "ppt_gradient_bg_v3_1770123504087.png")
ARCH_IMAGE = os.path.join(ARTIFACT_DIR, "project_architecture_v3_refined_1770123632340.png")
OUTPUT_FILE = os.path.join(ARTIFACT_DIR, "Student_Performance_Tracker_v3.pptx")

# Colors
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(173, 216, 230)
SOFT_GRAY = RGBColor(220, 220, 220)

def add_background(slide):
    if os.path.exists(BG_IMAGE):
        slide.shapes.add_picture(BG_IMAGE, 0, 0, Inches(10), Inches(5.625))

def set_text_style(shape, font_size=Pt(28), bold=False, color=WHITE, is_heading=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = 'Times New Roman'

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    title_box.text = "Student Assessment & Performance Tracker"
    set_text_style(title_box, font_size=Pt(48), bold=True, is_heading=True)
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    subtitle_box.text = "A Secure, Automated, and Data-Driven Educational Ecosystem"
    set_text_style(subtitle_box, font_size=Pt(20), color=LIGHT_BLUE)

    # 2. System Architecture (Visual)
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title.text = "System Architecture Overview"
    set_text_style(title, font_size=Pt(36), bold=True, is_heading=True)
    
    if os.path.exists(ARCH_IMAGE):
        # Center the diagram
        slide.shapes.add_picture(ARCH_IMAGE, Inches(2.2), Inches(1.0), height=Inches(3.8))
    
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.5))
    note_box.text = "Spring Boot Backend manages REST APIs, Security, and multi-channel Notifications (SMTP/Twilio)."
    set_text_style(note_box, font_size=Pt(14), color=SOFT_GRAY)
    for p in note_box.text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER

    # 3. Core Problem & Modern Solution
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title.text = "Problem vs Innovation"
    set_text_style(title, font_size=Pt(36), bold=True, is_heading=True)
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(3))
    left_box.text = ("Legacy Challenges:\n"
                     "• Fragmented data silos across departments.\n"
                     "• Manual reporting takes days to compile.\n"
                     "• Lack of real-time parent engagement.")
    set_text_style(left_box, font_size=Pt(18))
    
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.2), Inches(3))
    right_box.text = ("Modern Solution:\n"
                      "• Digital first, centralized SQL storage.\n"
                      "• Sub-second automated alert generation.\n"
                      "• Interactive, mobile-responsive portals.")
    set_text_style(right_box, font_size=Pt(18), color=LIGHT_BLUE)

    # 4. Functional Highlights
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title.text = "Functional Capabilities"
    set_text_style(title, font_size=Pt(36), bold=True, is_heading=True)
    
    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.5))
    content.text = ("• High-speed data processing for assessment marks and CGPA.\n"
                    "• Secure multi-role access (Admin, Faculty, Student, Parent).\n"
                    "• Automated bulk upload processing for CSV/Excel data.\n"
                    "• Dynamic trend visualization and performance heatmaps.\n"
                    "• Integrated microservices logic for reporting and analytics.")
    set_text_style(content, font_size=Pt(20))

    # 5. Technology & Security
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title.text = "Technology Stack & Security"
    set_text_style(title, font_size=Pt(36), bold=True, is_heading=True)
    
    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.5))
    content.text = ("• Backend: Java 17, Spring Boot 3 with Data JPA.\n"
                    "• Frontend: React 18, Vite for optimal performance.\n"
                    "• Notification: Twilio SMS Gateway & Java Mail (SMTP).\n"
                    "• Reporting: .NET Core dedicated analytics service.\n"
                    "• Security: JWT based authentication & role-based authorization.")
    set_text_style(content, font_size=Pt(20))

    # 6. Future Roadmap
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title.text = "Strategic Roadmap"
    set_text_style(title, font_size=Pt(36), bold=True, is_heading=True)
    
    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.5))
    content.text = ("• AI Integration: Predictive modeling for student performance.\n"
                    "• Smart Attendance: QR code based tracking system.\n"
                    "• Mobile Ecosystem: Dedicated Android/iOS parent applications.\n"
                    "• Cloud Migration: Enhanced scalability with AWS/Azure.\n"
                    "• Interactive Portals: Virtual Parent-Teacher meeting modules.")
    set_text_style(content, font_size=Pt(20))

    # 7. Conclusion
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    center_text = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    center_text.text = "Excellence Driven by Data\nTHANK YOU"
    set_text_style(center_text, font_size=Pt(44), bold=True, is_heading=True)
    for p in center_text.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_FILE)
    print(f"V3 Presentation saved to: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
